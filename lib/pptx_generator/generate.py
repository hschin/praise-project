#!/usr/bin/env python3
# Source: Noto Sans SC TTF bundled from Google Fonts (noto-cjk), converted to TTF via fonttools.
# Font is physically embedded into each .pptx output so it renders on Windows projectors
# without requiring font installation.

import base64
import json
import os
import sys
import tempfile
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.oxml import parse_xml
from lxml import etree

FONT_SIZE_MAP = {"small": 28, "medium": 36, "large": 48}
FONT_NAME = "Noto Sans SC"
FONT_PATH = os.path.join(os.path.dirname(__file__), "fonts", "NotoSansSC-Regular.ttf")
FONT_ZIP_NAME = "ppt/fonts/NotoSansSC-Regular.ttf"


def hex_to_rgb(hex_color):
    """Parse a hex color string (#rrggbb) into an RGBColor."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def set_ea_font(run, typeface):
    """Set the East Asian font element on a run's rPr so PowerPoint routes CJK glyphs correctly."""
    rPr = run._r.get_or_add_rPr()
    ea_elem = rPr.find(qn("a:ea"))
    if ea_elem is None:
        ea_elem = etree.SubElement(rPr, qn("a:ea"))
    ea_elem.set("typeface", typeface)


def set_para_spacing(para, space_before_pt=0, space_after_pt=0, line_spacing=1.0):
    """Set paragraph spacing via pPr XML — python-pptx has no high-level API for this."""
    pPr = para._p.get_or_add_pPr()
    if space_before_pt:
        spcBef = etree.SubElement(pPr, qn("a:spcBef"))
        spcPct = etree.SubElement(spcBef, qn("a:spcPts"))
        spcPct.set("val", str(int(space_before_pt * 100)))
    if space_after_pt:
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts.set("val", str(int(space_after_pt * 100)))
    if line_spacing != 1.0:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))


def add_title_slide(prs, song_title, song_artist, theme):
    """Add a song title slide before the lyric slides."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    bg_rgb = hex_to_rgb(theme["background_color"])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    # Background image (same as lyric slides)
    bg_b64 = theme.get("background_image_base64")
    if bg_b64:
        img_bytes = base64.b64decode(bg_b64)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
            tmp_img.write(img_bytes)
            tmp_img_path = tmp_img.name
        try:
            slide.shapes.add_picture(
                tmp_img_path,
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
        finally:
            if os.path.exists(tmp_img_path):
                os.unlink(tmp_img_path)

    text_rgb = hex_to_rgb(theme["text_color"])
    font_size_key = theme.get("font_size", "medium")
    title_pt = FONT_SIZE_MAP.get(font_size_key, FONT_SIZE_MAP["medium"])
    artist_pt = round(title_pt * 0.55)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    # Full-slide textbox with vertical centering
    txBox = slide.shapes.add_textbox(0, 0, slide_w, slide_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title line
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    set_para_spacing(para, line_spacing=1.2)
    run = para.add_run()
    run.text = song_title
    run.font.name = FONT_NAME
    run.font.size = Pt(title_pt)
    run.font.color.rgb = text_rgb
    run.font.bold = True
    set_ea_font(run, FONT_NAME)

    # Artist line (if present)
    if song_artist:
        para2 = tf.add_paragraph()
        para2.alignment = PP_ALIGN.CENTER
        set_para_spacing(para2, space_before_pt=title_pt * 0.3, line_spacing=1.2)
        run2 = para2.add_run()
        run2.text = song_artist
        run2.font.name = FONT_NAME
        run2.font.size = Pt(artist_pt)
        run2.font.color.rgb = text_rgb
        set_ea_font(run2, FONT_NAME)


def add_slide(prs, slide_data, theme):
    """Add a single content slide to the presentation."""
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Set background color
    bg_rgb = hex_to_rgb(theme["background_color"])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    # Handle optional background image
    bg_b64 = theme.get("background_image_base64")
    if bg_b64:
        img_bytes = base64.b64decode(bg_b64)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
            tmp_img.write(img_bytes)
            tmp_img_path = tmp_img.name
        try:
            slide.shapes.add_picture(
                tmp_img_path,
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
        finally:
            if os.path.exists(tmp_img_path):
                os.unlink(tmp_img_path)

    # Determine font sizes
    font_size_key = theme.get("font_size", "medium")
    chinese_pt = FONT_SIZE_MAP.get(font_size_key, FONT_SIZE_MAP["medium"])
    pinyin_pt = round(chinese_pt * 0.6)
    text_rgb = hex_to_rgb(theme["text_color"])

    # Full-slide textbox with vertical centering — matches preview flex centering
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    txBox = slide.shapes.add_textbox(
        left=int(slide_w * 0.08),
        top=0,
        width=int(slide_w * 0.84),
        height=slide_h,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    content = slide_data.get("content", "")
    pinyin = slide_data.get("pinyin", "")

    content_lines = content.split("\n") if content else []
    pinyin_lines = pinyin.split("\n") if pinyin else []

    first_para = True

    for i, chinese_line in enumerate(content_lines):
        p_line = pinyin_lines[i] if i < len(pinyin_lines) else ""

        # Pinyin line — tight spacing below (it sits just above the Chinese character)
        if p_line:
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            # No space before first pinyin; space before subsequent pairs = line gap
            set_para_spacing(para, space_before_pt=(chinese_pt * 0.5) if i > 0 else 0, line_spacing=1.0)
            run = para.add_run()
            run.text = p_line
            run.font.name = FONT_NAME
            run.font.size = Pt(pinyin_pt)
            run.font.color.rgb = text_rgb
            set_ea_font(run, FONT_NAME)

        # Chinese line — tight spacing below pinyin, normal line height
        if chinese_line:
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            set_para_spacing(para, space_before_pt=0, line_spacing=1.1)
            run = para.add_run()
            run.text = chinese_line
            run.font.name = FONT_NAME
            run.font.size = Pt(chinese_pt)
            run.font.color.rgb = text_rgb
            set_ea_font(run, FONT_NAME)


def embed_cjk_font(pptx_path):
    """Post-process the .pptx ZIP to embed the Noto Sans SC TTF binary.

    Steps:
    1. Read the saved .pptx into memory as a ZIP
    2. Add ppt/fonts/NotoSansSC-Regular.ttf from the bundled font file
    3. Update [Content_Types].xml to register the font part
    4. Add a relationship in ppt/_rels/presentation.xml.rels pointing to the font
    5. Write the modified ZIP back to pptx_path (via a temp file to avoid corruption)
    """
    tmp_path = pptx_path + ".tmp"
    with zipfile.ZipFile(pptx_path, "r") as zin, \
         zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:

        existing_names = set(zin.namelist())

        # Copy all existing entries, patching XML where needed
        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == "[Content_Types].xml":
                # Register the font part type if not already present
                font_override = (
                    '<Override PartName="/ppt/fonts/NotoSansSC-Regular.ttf" '
                    'ContentType="application/x-fontdata"/>'
                )
                if b"NotoSansSC" not in data:
                    data = data.replace(b"</Types>", font_override.encode() + b"</Types>")

            elif item.filename == "ppt/_rels/presentation.xml.rels":
                # Add a font relationship if not already present
                font_rel = (
                    '<Relationship Id="rFontNoto" '
                    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
                    'Target="fonts/NotoSansSC-Regular.ttf"/>'
                )
                if b"NotoSansSC" not in data:
                    data = data.replace(b"</Relationships>", font_rel.encode() + b"</Relationships>")

            zout.writestr(item, data)

        # Add the font binary (skip if somehow already present)
        if FONT_ZIP_NAME not in existing_names:
            zout.write(FONT_PATH, FONT_ZIP_NAME)

    # Replace original with modified
    os.replace(tmp_path, pptx_path)


def main():
    try:
        payload = json.loads(sys.stdin.read())
        deck = payload["deck"]
        output_path = payload["output_path"]
        theme = deck["theme"]

        # Validate output directory exists
        out_dir = os.path.dirname(os.path.abspath(output_path))
        if not os.path.isdir(out_dir):
            raise FileNotFoundError(f"Output directory does not exist: {out_dir}")

        # Create presentation (16:9 widescreen)
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Add slides for each song's lyric blocks
        for song in deck.get("songs", []):
            # Title slide first
            add_title_slide(prs, song.get("title", ""), song.get("artist", ""), theme)
            for slide_data in song.get("slides", []):
                add_slide(prs, slide_data, theme)

        prs.save(output_path)
        embed_cjk_font(output_path)

    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        sys.exit(1)

    sys.exit(0)


if __name__ == "__main__":
    main()
