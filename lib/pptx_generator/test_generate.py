#!/usr/bin/env python3
"""Tests for generate.py — PPTX generation script.

Run with: python3 -m pytest lib/pptx_generator/test_generate.py -v
or:        python3 lib/pptx_generator/test_generate.py
"""

import json
import os
import subprocess
import sys
import tempfile
import unittest
import zipfile

SCRIPT = os.path.join(os.path.dirname(__file__), "generate.py")
BASE_PAYLOAD = {
    "deck": {
        "title": "Test Service",
        "theme": {
            "background_color": "#1a1a2e",
            "text_color": "#ffffff",
            "font_size": "medium",
            "background_image_base64": None,
        },
        "songs": [
            {
                "title": "Test Song",
                "slides": [
                    {
                        "section_type": "verse",
                        "content": "你好世界",
                        "pinyin": "nǐ hǎo shì jiè",
                    }
                ],
            }
        ],
    },
    "output_path": None,  # filled per test
}


def run_script(payload):
    """Run generate.py with the given payload dict, return (returncode, stdout, stderr)."""
    result = subprocess.run(
        [sys.executable, SCRIPT],
        input=json.dumps(payload),
        capture_output=True,
        text=True,
    )
    return result.returncode, result.stdout, result.stderr


class TestGeneratePptx(unittest.TestCase):
    def make_payload(self, **overrides):
        import copy
        p = copy.deepcopy(BASE_PAYLOAD)
        p.update(overrides)
        return p

    def test_generates_one_slide_per_lyric_block(self):
        """One slide for each entry in songs[*].slides."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            out = f.name
        try:
            payload = self.make_payload()
            payload["output_path"] = out
            rc, _, stderr = run_script(payload)
            self.assertEqual(rc, 0, f"Script exited {rc}: {stderr}")
            from pptx import Presentation
            prs = Presentation(out)
            self.assertEqual(len(prs.slides), 1)
        finally:
            if os.path.exists(out):
                os.unlink(out)

    def test_cjk_font_embedded_in_zip(self):
        """ppt/fonts/NotoSansSC-Regular.ttf must be physically inside the .pptx archive."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            out = f.name
        try:
            payload = self.make_payload()
            payload["output_path"] = out
            rc, _, stderr = run_script(payload)
            self.assertEqual(rc, 0, f"Script exited {rc}: {stderr}")
            with zipfile.ZipFile(out) as z:
                self.assertIn(
                    "ppt/fonts/NotoSansSC-Regular.ttf",
                    z.namelist(),
                    "CJK font not embedded in .pptx ZIP",
                )
        finally:
            if os.path.exists(out):
                os.unlink(out)

    def test_multiple_songs_multiple_slides(self):
        """Two songs with 2 slides each → 4 slides total."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            out = f.name
        try:
            payload = self.make_payload()
            payload["output_path"] = out
            slide = {"section_type": "verse", "content": "测试内容", "pinyin": "cè shì nèi róng"}
            payload["deck"]["songs"] = [
                {"title": "Song A", "slides": [slide, slide]},
                {"title": "Song B", "slides": [slide, slide]},
            ]
            rc, _, stderr = run_script(payload)
            self.assertEqual(rc, 0, f"Script exited {rc}: {stderr}")
            from pptx import Presentation
            prs = Presentation(out)
            self.assertEqual(len(prs.slides), 4)
        finally:
            if os.path.exists(out):
                os.unlink(out)

    def test_exits_nonzero_on_invalid_output_path(self):
        """Script must exit non-zero when output directory does not exist."""
        payload = self.make_payload()
        payload["output_path"] = "/nonexistent/directory/file.pptx"
        rc, _, stderr = run_script(payload)
        self.assertNotEqual(rc, 0, "Expected non-zero exit for invalid output path")

    def test_pinyin_omitted_when_blank(self):
        """Script must succeed (exit 0) when pinyin field is blank."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            out = f.name
        try:
            payload = self.make_payload()
            payload["output_path"] = out
            payload["deck"]["songs"][0]["slides"][0]["pinyin"] = ""
            rc, _, stderr = run_script(payload)
            self.assertEqual(rc, 0, f"Script exited {rc}: {stderr}")
        finally:
            if os.path.exists(out):
                os.unlink(out)

    def test_exits_zero_on_success(self):
        """Script exits 0 on success."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            out = f.name
        try:
            payload = self.make_payload()
            payload["output_path"] = out
            rc, _, _ = run_script(payload)
            self.assertEqual(rc, 0)
        finally:
            if os.path.exists(out):
                os.unlink(out)


if __name__ == "__main__":
    unittest.main()
